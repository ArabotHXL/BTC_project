#!/usr/bin/env python3
"""
Generate Chinese Product Introduction PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_chinese_ppt():
    """Create Chinese product introduction PowerPoint"""
    print("📊 Creating Chinese Product Introduction PPT...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)  # Dark background
    
    # Main title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "HashInsight 比特币挖矿管理平台"
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "专业 • 安全 • 智能"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = RGBColor(243, 156, 18)  # Gold accent
    
    # Slide 2: Three-column layout
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "核心优势"
    title_p2 = title_frame2.paragraphs[0]
    title_p2.alignment = PP_ALIGN.CENTER
    title_p2.font.size = Pt(36)
    title_p2.font.bold = True
    title_p2.font.color.rgb = RGBColor(44, 62, 80)
    
    # Decorative line
    line = slide2.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(243, 156, 18)
    line.line.color.rgb = RGBColor(243, 156, 18)
    
    # Column 1: 核心功能
    col1_x = Inches(0.5)
    col1_y = Inches(1.5)
    col_width = Inches(2.8)
    
    # Column 1 header box with dark background
    col1_header = slide2.shapes.add_textbox(col1_x, col1_y, col_width, Inches(0.6))
    col1_header_frame = col1_header.text_frame
    col1_header_frame.text = "核心功能"
    col1_header_p = col1_header_frame.paragraphs[0]
    col1_header_p.alignment = PP_ALIGN.CENTER
    col1_header_p.font.size = Pt(22)
    col1_header_p.font.bold = True
    col1_header_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add background shape for header
    header_bg1 = slide2.shapes.add_shape(1, col1_x, col1_y, col_width, Inches(0.6))
    header_bg1.fill.solid()
    header_bg1.fill.fore_color.rgb = RGBColor(52, 73, 94)
    header_bg1.line.color.rgb = RGBColor(52, 73, 94)
    slide2.shapes._spTree.remove(header_bg1._element)
    slide2.shapes._spTree.insert(2, header_bg1._element)
    
    # Column 1 content
    col1_content = slide2.shapes.add_textbox(col1_x, col1_y + Inches(0.7), col_width, Inches(4.5))
    col1_content_frame = col1_content.text_frame
    col1_content_frame.word_wrap = True
    col1_content_frame.margin_left = Pt(10)
    col1_content_frame.margin_right = Pt(10)
    
    items1 = [
        "• 智能挖矿计算器",
        "  （双算法验证，",
        "   17+矿机型号）",
        "",
        "• 技术分析平台",
        "  （10个信号模块）",
        "",
        "• 资金管理系统",
        "  （BTC库存+",
        "   卖出策略）",
        "",
        "• CRM客户管理",
        "  （全生命周期）",
        "",
        "• 批量计算",
        "  （5000+矿机）"
    ]
    
    for item in items1:
        p = col1_content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(2)
        if item.startswith("•"):
            p.font.bold = True
            p.font.color.rgb = RGBColor(52, 73, 94)
        else:
            p.font.color.rgb = RGBColor(80, 80, 80)
    
    # Column 2: 安全保障
    col2_x = Inches(3.6)
    col2_y = Inches(1.5)
    
    # Column 2 header
    col2_header = slide2.shapes.add_textbox(col2_x, col2_y, col_width, Inches(0.6))
    col2_header_frame = col2_header.text_frame
    col2_header_frame.text = "安全保障"
    col2_header_p = col2_header_frame.paragraphs[0]
    col2_header_p.alignment = PP_ALIGN.CENTER
    col2_header_p.font.size = Pt(22)
    col2_header_p.font.bold = True
    col2_header_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add background shape for header
    header_bg2 = slide2.shapes.add_shape(1, col2_x, col2_y, col_width, Inches(0.6))
    header_bg2.fill.solid()
    header_bg2.fill.fore_color.rgb = RGBColor(41, 128, 185)
    header_bg2.line.color.rgb = RGBColor(41, 128, 185)
    slide2.shapes._spTree.remove(header_bg2._element)
    slide2.shapes._spTree.insert(2, header_bg2._element)
    
    # Column 2 content
    col2_content = slide2.shapes.add_textbox(col2_x, col2_y + Inches(0.7), col_width, Inches(4.5))
    col2_content_frame = col2_content.text_frame
    col2_content_frame.word_wrap = True
    col2_content_frame.margin_left = Pt(10)
    col2_content_frame.margin_right = Pt(10)
    
    items2 = [
        "• 企业级加密",
        "  （KMS密钥管理）",
        "",
        "• mTLS身份认证",
        "",
        "• WireGuard VPN",
        "",
        "• SOC 2合规",
        "  准备中",
        "",
        "• 99.95% SLA保证"
    ]
    
    for item in items2:
        p = col2_content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(2)
        if item.startswith("•"):
            p.font.bold = True
            p.font.color.rgb = RGBColor(41, 128, 185)
        else:
            p.font.color.rgb = RGBColor(80, 80, 80)
    
    # Column 3: 适用场景
    col3_x = Inches(6.7)
    col3_y = Inches(1.5)
    
    # Column 3 header
    col3_header = slide2.shapes.add_textbox(col3_x, col3_y, col_width, Inches(0.6))
    col3_header_frame = col3_header.text_frame
    col3_header_frame.text = "适用场景"
    col3_header_p = col3_header_frame.paragraphs[0]
    col3_header_p.alignment = PP_ALIGN.CENTER
    col3_header_p.font.size = Pt(22)
    col3_header_p.font.bold = True
    col3_header_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add background shape for header
    header_bg3 = slide2.shapes.add_shape(1, col3_x, col3_y, col_width, Inches(0.6))
    header_bg3.fill.solid()
    header_bg3.fill.fore_color.rgb = RGBColor(243, 156, 18)
    header_bg3.line.color.rgb = RGBColor(243, 156, 18)
    slide2.shapes._spTree.remove(header_bg3._element)
    slide2.shapes._spTree.insert(2, header_bg3._element)
    
    # Column 3 content
    col3_content = slide2.shapes.add_textbox(col3_x, col3_y + Inches(0.7), col_width, Inches(4.5))
    col3_content_frame = col3_content.text_frame
    col3_content_frame.word_wrap = True
    col3_content_frame.margin_left = Pt(10)
    col3_content_frame.margin_right = Pt(10)
    
    items3 = [
        "• 大型矿场",
        "  （100MW+，",
        "   15000+矿机）",
        "",
        "• 托管服务商",
        "  （300+客户管理）",
        "",
        "• 机构投资者",
        "  （专业资金管理）",
        "",
        "• 矿机经销商",
        "  （批量报价）"
    ]
    
    for item in items3:
        p = col3_content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.space_after = Pt(2)
        if item.startswith("•"):
            p.font.bold = True
            p.font.color.rgb = RGBColor(243, 156, 18)
        else:
            p.font.color.rgb = RGBColor(80, 80, 80)
    
    # Save the presentation
    output_file = "HashInsight_产品介绍_简版_CN.pptx"
    prs.save(output_file)
    print(f"✅ Chinese PowerPoint created: {output_file}")

def main():
    print("\n" + "="*70)
    print("生成中文产品介绍PPT")
    print("="*70 + "\n")
    
    create_chinese_ppt()
    
    print("\n" + "="*70)
    print("✅ 中文PPT生成成功！")
    print("="*70)
    print("\n文件: HashInsight_产品介绍_简版_CN.pptx")

if __name__ == '__main__':
    main()
