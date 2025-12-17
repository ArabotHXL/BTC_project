#!/usr/bin/env python3
"""
BTC Mining Calculator 系统演示PPT生成器
使用python-pptx库创建专业的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """创建BTC Mining Calculator演示PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色主题
    BG_DARK = RGBColor(26, 29, 46)  # #1a1d2e
    GOLD = RGBColor(247, 147, 26)    # #f7931a
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(200, 200, 200)
    
    # 幻灯片1: 封面
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_background(slide1, BG_DARK)
    
    # 标题
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "BTC Mining Calculator"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = GOLD
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Enterprise-Grade Bitcoin Mining Analytics Platform"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 中文副标题
    subtitle_cn_box = slide1.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.6))
    subtitle_cn_frame = subtitle_cn_box.text_frame
    subtitle_cn_frame.text = "企业级比特币挖矿分析平台"
    subtitle_cn_frame.paragraphs[0].font.size = Pt(20)
    subtitle_cn_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    subtitle_cn_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 幻灯片2: 系统概述
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, BG_DARK)
    add_title(slide2, "System Overview | 系统概述", GOLD, WHITE)
    
    content2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf2 = content2.text_frame
    tf2.word_wrap = True
    
    add_bullet_point(tf2, "🎯 Real-time Bitcoin mining profitability analysis", WHITE, 20)
    add_bullet_point(tf2, "   实时比特币挖矿盈利能力分析", LIGHT_GRAY, 16)
    add_bullet_point(tf2, "", WHITE, 10)
    
    add_bullet_point(tf2, "📊 Supports 17+ ASIC miner models", WHITE, 20)
    add_bullet_point(tf2, "   支持17种以上ASIC矿机型号", LIGHT_GRAY, 16)
    add_bullet_point(tf2, "", WHITE, 10)
    
    add_bullet_point(tf2, "🌐 Complete bilingual support (English/Chinese)", WHITE, 20)
    add_bullet_point(tf2, "   完整中英文双语支持", LIGHT_GRAY, 16)
    add_bullet_point(tf2, "", WHITE, 10)
    
    add_bullet_point(tf2, "⛓️ Web3 integration with blockchain transparency", WHITE, 20)
    add_bullet_point(tf2, "   Web3集成与区块链透明度", LIGHT_GRAY, 16)
    
    # 幻灯片3: 核心功能模块
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, BG_DARK)
    add_title(slide3, "Core Modules | 核心功能", GOLD, WHITE)
    
    # 左列
    left_col = slide3.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    tf_left = left_col.text_frame
    tf_left.word_wrap = True
    
    add_module_item(tf_left, "1. Mining Calculator", "挖矿计算器", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• Dual-algorithm profitability analysis", "  双算法盈利能力分析", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• ROI & breakeven calculations", "  投资回报率与盈亏平衡计算", WHITE, LIGHT_GRAY)
    add_bullet_point(tf_left, "", WHITE, 10)
    
    add_module_item(tf_left, "2. CRM System", "CRM客户管理", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• 60+ API endpoints", "  60+个API接口", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• Lead & sales funnel management", "  潜在客户与销售漏斗管理", WHITE, LIGHT_GRAY)
    add_bullet_point(tf_left, "", WHITE, 10)
    
    add_module_item(tf_left, "3. System Monitoring", "系统监控", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• Real-time health checks", "  实时健康检查", WHITE, LIGHT_GRAY)
    add_module_item(tf_left, "• Performance metrics & alerts", "  性能指标与警报", WHITE, LIGHT_GRAY)
    
    # 右列
    right_col = slide3.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(5))
    tf_right = right_col.text_frame
    tf_right.word_wrap = True
    
    add_module_item(tf_right, "4. Technical Analysis", "技术分析", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• RSI, MACD, Bollinger Bands", "  RSI、MACD、布林带", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• Historical price analysis", "  历史价格分析", WHITE, LIGHT_GRAY)
    add_bullet_point(tf_right, "", WHITE, 10)
    
    add_module_item(tf_right, "5. Hosting Services", "托管服务", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• Miner management & monitoring", "  矿机管理与监控", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• Real-time telemetry data", "  实时遥测数据", WHITE, LIGHT_GRAY)
    add_bullet_point(tf_right, "", WHITE, 10)
    
    add_module_item(tf_right, "6. Web3 Integration", "Web3集成", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• Blockchain verification", "  区块链验证", WHITE, LIGHT_GRAY)
    add_module_item(tf_right, "• SLA NFT certificates", "  SLA NFT证书", WHITE, LIGHT_GRAY)
    
    # 幻灯片4: 技术架构
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, BG_DARK)
    add_title(slide4, "Technical Architecture | 技术架构", GOLD, WHITE)
    
    content4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(5))
    tf4 = content4.text_frame
    tf4.word_wrap = True
    
    add_bullet_point(tf4, "Backend | 后端", GOLD, 22, True)
    add_bullet_point(tf4, "• Flask + SQLAlchemy + PostgreSQL", WHITE, 18)
    add_bullet_point(tf4, "• Redis for caching & task queuing", LIGHT_GRAY, 16)
    add_bullet_point(tf4, "• RESTful API architecture", LIGHT_GRAY, 16)
    add_bullet_point(tf4, "", WHITE, 10)
    
    add_bullet_point(tf4, "Frontend | 前端", GOLD, 22, True)
    add_bullet_point(tf4, "• Jinja2 + Bootstrap 5", WHITE, 18)
    add_bullet_point(tf4, "• Chart.js for data visualization", LIGHT_GRAY, 16)
    add_bullet_point(tf4, "• Responsive mobile-first design", LIGHT_GRAY, 16)
    add_bullet_point(tf4, "", WHITE, 10)
    
    add_bullet_point(tf4, "Data Integration | 数据集成", GOLD, 22, True)
    add_bullet_point(tf4, "• Multi-exchange API integration (Binance, OKX, Deribit, Bybit)", WHITE, 16)
    add_bullet_point(tf4, "• Intelligent fallback & SWR caching", LIGHT_GRAY, 16)
    
    # 幻灯片5: Web3 & 区块链集成
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, BG_DARK)
    add_title(slide5, "Web3 & Blockchain Integration | Web3与区块链集成", GOLD, WHITE)
    
    content5 = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf5 = content5.text_frame
    tf5.word_wrap = True
    
    add_bullet_point(tf5, "🔗 Blockchain Transparency", GOLD, 24, True)
    add_bullet_point(tf5, "• Data verification on Base L2 network", WHITE, 18)
    add_bullet_point(tf5, "• IPFS distributed storage via Pinata", LIGHT_GRAY, 16)
    add_bullet_point(tf5, "• Immutable audit trail", LIGHT_GRAY, 16)
    add_bullet_point(tf5, "", WHITE, 10)
    
    add_bullet_point(tf5, "🎫 SLA NFT Certificates", GOLD, 24, True)
    add_bullet_point(tf5, "• Automated SLA compliance minting", WHITE, 18)
    add_bullet_point(tf5, "• On-chain performance guarantees", LIGHT_GRAY, 16)
    add_bullet_point(tf5, "• Transparent hosting accountability", LIGHT_GRAY, 16)
    add_bullet_point(tf5, "", WHITE, 10)
    
    add_bullet_point(tf5, "🔐 Secure Configuration Wizard", GOLD, 24, True)
    add_bullet_point(tf5, "• MetaMask wallet integration", WHITE, 18)
    add_bullet_point(tf5, "• Environment-based secrets management", LIGHT_GRAY, 16)
    
    # 幻灯片6: 托管服务功能
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6, BG_DARK)
    add_title(slide6, "Hosting Services | 托管服务", GOLD, WHITE)
    
    content6 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf6 = content6.text_frame
    tf6.word_wrap = True
    
    add_bullet_point(tf6, "⛏️ Miner Management | 矿机管理", GOLD, 22, True)
    add_bullet_point(tf6, "• Single & batch miner registration", WHITE, 18)
    add_bullet_point(tf6, "• Approval workflow for customer submissions", LIGHT_GRAY, 16)
    add_bullet_point(tf6, "• Serial number tracking & inventory", LIGHT_GRAY, 16)
    add_bullet_point(tf6, "", WHITE, 10)
    
    add_bullet_point(tf6, "📊 Real-Time Monitoring | 实时监控", GOLD, 22, True)
    add_bullet_point(tf6, "• Hashrate & power consumption tracking", WHITE, 18)
    add_bullet_point(tf6, "• Temperature & fan speed telemetry", LIGHT_GRAY, 16)
    add_bullet_point(tf6, "• Mining pool statistics (shares, workers)", LIGHT_GRAY, 16)
    add_bullet_point(tf6, "", WHITE, 10)
    
    add_bullet_point(tf6, "🏢 Site Management | 站点管理", GOLD, 22, True)
    add_bullet_point(tf6, "• Multi-site capacity tracking", WHITE, 18)
    add_bullet_point(tf6, "• Incident & ticket management", LIGHT_GRAY, 16)
    
    # 幻灯片7: 数据与分析
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide7, BG_DARK)
    add_title(slide7, "Data & Analytics | 数据与分析", GOLD, WHITE)
    
    content7 = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf7 = content7.text_frame
    tf7.word_wrap = True
    
    add_bullet_point(tf7, "📈 Market Data Collection | 市场数据收集", GOLD, 22, True)
    add_bullet_point(tf7, "• 4 major exchanges with 100% data completeness", WHITE, 18)
    add_bullet_point(tf7, "• Network hashrate & difficulty tracking", LIGHT_GRAY, 16)
    add_bullet_point(tf7, "• Automated every 15 minutes", LIGHT_GRAY, 16)
    add_bullet_point(tf7, "", WHITE, 10)
    
    add_bullet_point(tf7, "🧠 Intelligence Layer | 智能分析层", GOLD, 22, True)
    add_bullet_point(tf7, "• Predictive analytics with XGBoost", WHITE, 18)
    add_bullet_point(tf7, "• Anomaly detection algorithms", LIGHT_GRAY, 16)
    add_bullet_point(tf7, "• Power optimization recommendations", LIGHT_GRAY, 16)
    add_bullet_point(tf7, "", WHITE, 10)
    
    add_bullet_point(tf7, "📊 Technical Indicators | 技术指标", GOLD, 22, True)
    add_bullet_point(tf7, "• RSI, MACD, SMA, EMA, Bollinger Bands", WHITE, 18)
    
    # 幻灯片8: 安全特性
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide8, BG_DARK)
    add_title(slide8, "Security Features | 安全特性", GOLD, WHITE)
    
    content8 = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    tf8 = content8.text_frame
    tf8.word_wrap = True
    
    add_bullet_point(tf8, "🔐 Authentication & Authorization | 认证与授权", GOLD, 22, True)
    add_bullet_point(tf8, "• Custom email-based authentication", WHITE, 18)
    add_bullet_point(tf8, "• Role-based access control (RBAC)", LIGHT_GRAY, 16)
    add_bullet_point(tf8, "• Session management with secure tokens", LIGHT_GRAY, 16)
    add_bullet_point(tf8, "", WHITE, 10)
    
    add_bullet_point(tf8, "🛡️ Data Protection | 数据保护", GOLD, 22, True)
    add_bullet_point(tf8, "• AES encryption for sensitive data", WHITE, 18)
    add_bullet_point(tf8, "• Secure password hashing (Werkzeug)", LIGHT_GRAY, 16)
    add_bullet_point(tf8, "• Environment-based secrets (Replit Secrets)", LIGHT_GRAY, 16)
    add_bullet_point(tf8, "", WHITE, 10)
    
    add_bullet_point(tf8, "⛓️ Blockchain Security | 区块链安全", GOLD, 22, True)
    add_bullet_point(tf8, "• MetaMask integration (no key sharing)", WHITE, 18)
    add_bullet_point(tf8, "• Testnet-first development approach", LIGHT_GRAY, 16)
    
    # 幻灯片9: 系统亮点
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide9, BG_DARK)
    add_title(slide9, "Key Highlights | 系统亮点", GOLD, WHITE)
    
    # 创建3列布局
    col1 = slide9.shapes.add_textbox(Inches(0.3), Inches(2), Inches(3.1), Inches(5))
    tf_col1 = col1.text_frame
    tf_col1.word_wrap = True
    
    add_bullet_point(tf_col1, "📊 Scale", GOLD, 20, True)
    add_bullet_point(tf_col1, "• 17+ miner models", WHITE, 16)
    add_bullet_point(tf_col1, "• 60+ API endpoints", WHITE, 16)
    add_bullet_point(tf_col1, "• 56+ KPI cards", WHITE, 16)
    add_bullet_point(tf_col1, "• 42+ visualizations", WHITE, 16)
    
    col2 = slide9.shapes.add_textbox(Inches(3.4), Inches(2), Inches(3.1), Inches(5))
    tf_col2 = col2.text_frame
    tf_col2.word_wrap = True
    
    add_bullet_point(tf_col2, "⚡ Performance", GOLD, 20, True)
    add_bullet_point(tf_col2, "• Real-time updates", WHITE, 16)
    add_bullet_point(tf_col2, "• Redis caching", WHITE, 16)
    add_bullet_point(tf_col2, "• Optimized queries", WHITE, 16)
    add_bullet_point(tf_col2, "• Auto-scaling ready", WHITE, 16)
    
    col3 = slide9.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3.1), Inches(5))
    tf_col3 = col3.text_frame
    tf_col3.word_wrap = True
    
    add_bullet_point(tf_col3, "🌐 Integration", GOLD, 20, True)
    add_bullet_point(tf_col3, "• 4 exchanges", WHITE, 16)
    add_bullet_point(tf_col3, "• Web3 blockchain", WHITE, 16)
    add_bullet_point(tf_col3, "• IPFS storage", WHITE, 16)
    add_bullet_point(tf_col3, "• NFT minting", WHITE, 16)
    
    # 幻灯片10: 总结与展望
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide10, BG_DARK)
    add_title(slide10, "Summary | 总结", GOLD, WHITE)
    
    content10 = slide10.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
    tf10 = content10.text_frame
    tf10.word_wrap = True
    
    add_bullet_point(tf10, "✅ Comprehensive Bitcoin mining analytics platform", WHITE, 20)
    add_bullet_point(tf10, "   全面的比特币挖矿分析平台", LIGHT_GRAY, 18)
    add_bullet_point(tf10, "", WHITE, 10)
    
    add_bullet_point(tf10, "✅ Enterprise-ready with robust architecture", WHITE, 20)
    add_bullet_point(tf10, "   企业级架构，稳定可靠", LIGHT_GRAY, 18)
    add_bullet_point(tf10, "", WHITE, 10)
    
    add_bullet_point(tf10, "✅ Web3-enabled transparency and accountability", WHITE, 20)
    add_bullet_point(tf10, "   Web3驱动的透明度与问责制", LIGHT_GRAY, 18)
    add_bullet_point(tf10, "", WHITE, 10)
    
    add_bullet_point(tf10, "✅ Real-time monitoring and intelligent insights", WHITE, 20)
    add_bullet_point(tf10, "   实时监控与智能洞察", LIGHT_GRAY, 18)
    
    # 幻灯片11: Thank You
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide11, BG_DARK)
    
    thank_you_box = slide11.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You\n谢谢"
    for paragraph in thank_you_frame.paragraphs:
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.color.rgb = GOLD
        paragraph.alignment = PP_ALIGN.CENTER
    
    contact_box = slide11.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "BTC Mining Calculator Platform\nEnterprise-Grade Mining Analytics"
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = "BTC_Mining_Calculator_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PPT生成成功: {output_path}")
    return output_path

def set_slide_background(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, color, subtitle_color=None):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_bullet_point(text_frame, text, color, size, bold=False):
    """添加项目符号点"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = 0

def add_module_item(text_frame, main_text, sub_text, main_color, sub_color):
    """添加模块项（包含主文本和子文本）"""
    p1 = text_frame.add_paragraph()
    p1.text = main_text
    p1.font.size = Pt(16)
    p1.font.color.rgb = main_color
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = sub_text
    p2.font.size = Pt(14)
    p2.font.color.rgb = sub_color

if __name__ == "__main__":
    create_presentation()
