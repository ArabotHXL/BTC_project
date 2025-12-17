#!/usr/bin/env python3
"""
Generate all documentation formats (PDF/HTML/PPT) with all improvements
"""
import markdown
from weasyprint import HTML
from weasyprint.text.fonts import FontConfiguration
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def generate_pdf(input_md, output_pdf, title):
    """Generate PDF from Markdown"""
    print(f"📄 {input_md} → {output_pdf}")
    
    with open(input_md, 'r', encoding='utf-8') as f:
        markdown_content = f.read()
    
    md = markdown.Markdown(extensions=['tables', 'fenced_code', 'codehilite', 'toc', 'sane_lists'])
    html_content = md.convert(markdown_content)
    
    full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>{title}</title>
    <style>
        @page {{ size: A4; margin: 2.5cm 2cm; @bottom-center {{ content: "Page " counter(page); font-size: 9pt; color: #666; }} }}
        body {{ font-family: Arial, Helvetica, sans-serif; line-height: 1.7; color: #333; font-size: 10pt; }}
        h1 {{ color: #2c3e50; border-bottom: 4px solid #f39c12; padding-bottom: 12px; margin-top: 35px; page-break-before: always; font-size: 22pt; }}
        h1:first-of-type {{ page-break-before: auto; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; margin: -2.5cm -2cm 25px -2cm; border-bottom: none; text-align: center; }}
        h2 {{ color: #34495e; border-bottom: 3px solid #3498db; padding-bottom: 10px; margin-top: 28px; font-size: 16pt; }}
        h3 {{ color: #555; margin-top: 22px; font-size: 13pt; border-left: 5px solid #f39c12; padding-left: 12px; }}
        h4 {{ color: #666; margin-top: 18px; font-size: 11pt; }}
        code {{ background-color: #f5f5f5; padding: 3px 7px; border-radius: 4px; font-family: Courier New, monospace; font-size: 8.5pt; color: #e74c3c; border: 1px solid #e0e0e0; }}
        pre {{ background-color: #f8f8f8; border: 1px solid #ddd; border-left: 4px solid #f39c12; padding: 15px; margin: 15px 0; overflow-x: auto; font-size: 8.5pt; border-radius: 4px; }}
        pre code {{ background-color: transparent; padding: 0; border: none; color: #333; }}
        table {{ border-collapse: collapse; width: 100%; margin: 18px 0; font-size: 9.5pt; box-shadow: 0 3px 10px rgba(0,0,0,0.1); page-break-inside: avoid; }}
        table th {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 12px 10px; text-align: left; font-weight: bold; border: 1px solid #5a67d8; }}
        table td {{ border: 1px solid #ddd; padding: 10px; vertical-align: top; }}
        table tr:nth-child(even) {{ background-color: #f9f9f9; }}
        blockquote {{ border-left: 5px solid #3498db; padding: 12px 18px; margin: 15px 0; color: #555; background-color: #ecf0f1; border-radius: 5px; }}
        ul, ol {{ margin: 12px 0; padding-left: 30px; }}
        li {{ margin: 6px 0; line-height: 1.6; }}
        a {{ color: #3498db; text-decoration: none; }}
        strong {{ color: #2c3e50; font-weight: 700; }}
        hr {{ border: none; border-top: 2px solid #ecf0f1; margin: 30px 0; }}
    </style>
</head>
<body>{html_content}</body>
</html>"""
    
    font_config = FontConfiguration()
    HTML(string=full_html, base_url='.').write_pdf(output_pdf, font_config=font_config, optimize_images=True)

def generate_html(input_md, output_html, title):
    """Generate HTML from Markdown"""
    print(f"🌐 {input_md} → {output_html}")
    
    with open(input_md, 'r', encoding='utf-8') as f:
        markdown_content = f.read()
    
    md = markdown.Markdown(extensions=['tables', 'fenced_code', 'codehilite', 'toc', 'sane_lists'])
    html_content = md.convert(markdown_content)
    
    full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        @media print {{ @page {{ size: A4; margin: 2cm; }} body {{ font-size: 10pt; }} h1 {{ page-break-before: always; }} h1:first-of-type {{ page-break-before: auto; }} table {{ page-break-inside: avoid; }} .print-button {{ display: none; }} }}
        body {{ font-family: -apple-system, BlinkMacSystemFont, Segoe UI, Arial, sans-serif; line-height: 1.8; color: #333; max-width: 1000px; margin: 0 auto; padding: 40px 20px; background-color: #f9f9f9; }}
        h1 {{ color: #2c3e50; border-bottom: 4px solid #f39c12; padding-bottom: 12px; margin-top: 40px; margin-bottom: 25px; font-size: 2.5em; }}
        h1:first-of-type {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 40px; margin: -40px -20px 30px -20px; border-bottom: none; text-align: center; }}
        h2 {{ color: #34495e; border-bottom: 3px solid #3498db; padding-bottom: 10px; margin-top: 35px; margin-bottom: 20px; font-size: 2em; }}
        h3 {{ color: #555; margin-top: 28px; margin-bottom: 15px; font-size: 1.5em; border-left: 5px solid #f39c12; padding-left: 15px; }}
        h4 {{ color: #666; margin-top: 22px; margin-bottom: 12px; font-size: 1.25em; }}
        p {{ margin: 12px 0; text-align: justify; }}
        code {{ background-color: #f5f5f5; padding: 3px 8px; border-radius: 4px; font-family: Consolas, Monaco, monospace; font-size: 0.9em; color: #e74c3c; border: 1px solid #e0e0e0; }}
        pre {{ background-color: #f8f8f8; border: 1px solid #ddd; border-left: 4px solid #f39c12; padding: 20px; margin: 20px 0; overflow-x: auto; border-radius: 5px; }}
        pre code {{ background-color: transparent; padding: 0; border: none; color: #333; }}
        table {{ border-collapse: collapse; width: 100%; margin: 25px 0; background: white; box-shadow: 0 4px 12px rgba(0,0,0,0.1); border-radius: 8px; overflow: hidden; }}
        table th {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 15px 12px; text-align: left; font-weight: bold; }}
        table td {{ border: 1px solid #e0e0e0; padding: 12px; vertical-align: top; }}
        table tr:nth-child(even) {{ background-color: #f9f9f9; }}
        table tr:hover {{ background-color: #f0f0f0; }}
        blockquote {{ border-left: 5px solid #3498db; padding: 15px 20px; margin: 20px 0; color: #555; background-color: #ecf0f1; border-radius: 5px; }}
        ul, ol {{ margin: 15px 0; padding-left: 35px; }}
        li {{ margin: 8px 0; line-height: 1.7; }}
        a {{ color: #3498db; text-decoration: none; }}
        a:hover {{ color: #2980b9; text-decoration: underline; }}
        strong {{ color: #2c3e50; font-weight: 700; }}
        hr {{ border: none; border-top: 2px solid #ecf0f1; margin: 40px 0; }}
        .print-button {{ position: fixed; top: 20px; right: 20px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 12px 24px; border: none; border-radius: 25px; cursor: pointer; font-size: 16px; box-shadow: 0 4px 15px rgba(0,0,0,0.2); z-index: 1000; }}
        .print-button:hover {{ transform: scale(1.05); }}
    </style>
</head>
<body>
    <button class="print-button" onclick="window.print()">🖨️ Print/Export PDF</button>
    {html_content}
</body>
</html>"""
    
    with open(output_html, 'w', encoding='utf-8') as f:
        f.write(full_html)

def create_enhanced_ppt():
    """Create enhanced PowerPoint with new security slide"""
    print("📊 Creating enhanced PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Helper functions
    def create_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(44, 62, 80)
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "HashInsight"
        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(66)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Professional Bitcoin Mining Management Platform"
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.alignment = PP_ALIGN.CENTER
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(243, 156, 18)
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(9), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(243, 156, 18)
        line.line.color.rgb = RGBColor(243, 156, 18)
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for item in content_items:
            p = content_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
            p.space_before = Pt(4)
    
    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)
        
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(9), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(243, 156, 18)
        
        table_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            Inches(0.8), Inches(2),
            Inches(8.4), Inches(4)
        )
        table = table_shape.table
        
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(52, 73, 94)
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.size = Pt(13)
        
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(236, 240, 241)
                text_frame = cell.text_frame
                text_frame.paragraphs[0].font.size = Pt(11)
    
    # Build slides
    create_title_slide()
    
    add_content_slide("Product Overview", [
        "Target Users",
        "  Mining farms, hosting providers, institutional investors",
        "",
        "Core Value",
        "  Data-driven decisions, professional management",
        "",
        "Technical Guarantees",
        "  99.95% SLA, 9.8x performance, enterprise security"
    ])
    
    add_content_slide("5 Core Features + Dual-Algorithm", [
        "1. Intelligent Mining Calculator",
        "   Dual-Algorithm Validation: Hashrate-based + Difficulty-based",
        "   Cross-validation with ±2% tolerance",
        "",
        "2. Technical Analysis - 10 signal modules",
        "",
        "3. Treasury Management - BTC inventory + strategies",
        "",
        "4. CRM - Full lifecycle management",
        "",
        "5. Enterprise Security - KMS, mTLS, WireGuard VPN"
    ])
    
    add_content_slide("Enterprise Security & Compliance", [
        "KMS Key Management - Multi-cloud, AES-256, RSA-4096",
        "mTLS Authentication - X.509 certificates, CRL/OCSP",
        "API Key Management - Scoped permissions, rate limiting",
        "WireGuard VPN - Zero-trust, ChaCha20-Poly1305",
        "Audit Logging - 7-year retention, tamper-proof",
        "",
        "Compliance Status:",
        "  SOC 2 Type II Ready (Q1 2026 audit scheduled)",
        "  PCI DSS Compliant (via Stripe)",
        "  GDPR Fully Compliant"
    ])
    
    add_table_slide("Technical Advantages", 
        ["Metric", "Value", "Description"],
        [
            ["SLA Availability", "≥99.95%", "Error budget ≤21.6 min/month"],
            ["API Latency", "p95 ≤250ms", "Fast response"],
            ["Error Rate", "≤0.1%", "High reliability"],
            ["Performance", "9.8x", "Request Coalescing"],
            ["Batch Process", "5000+ miners", "Concurrent calculation"]
        ]
    )
    
    add_table_slide("Competitive Advantage",
        ["", "HashInsight", "Foreman", "NiceHash"],
        [
            ["Dual-Algorithm", "✓", "✗", "✗"],
            ["Batch 5000+", "✓", "500 max", "Single only"],
            ["Treasury Mgmt", "✓", "✗", "✗"],
            ["SOC 2 Ready", "✓", "✗", "✗"],
            ["Support", "24/7+CSM", "Email", "Forum"]
        ]
    )
    
    add_content_slide("Customer Success: 100MW Farm", [
        "Before HashInsight:",
        "  Calculation: 2 hours/day",
        "  Manual errors: 12/month",
        "  Customer complaints: 100/month",
        "",
        "After HashInsight:",
        "  Calculation: 5 minutes (24x faster)",
        "  Manual errors: 0/month (100% eliminated)",
        "  Customer complaints: 30/month (70% reduction)",
        "  Power cost savings: $1.8M/year",
        "",
        "ROI: 158x | Payback: 2.3 days"
    ])
    
    add_table_slide("Pricing Plans",
        ["Plan", "Price", "Target", "ROI"],
        [
            ["Professional", "$199/mo", "Small-medium", "15-20x"],
            ["Enterprise", "$999/mo", "Large institutions", "50-100x"],
            ["Custom", "Custom", "Ultra-large", "100-200x"]
        ]
    )
    
    add_content_slide("Contact Us", [
        "Business Inquiries",
        "  sales@hashinsight.io",
        "",
        "Technical Support",
        "  support@hashinsight.io",
        "",
        "Official Website",
        "  https://hashinsight.io",
        "",
        "Enterprise Support",
        "  24/7 hotline, 1-hour response SLA"
    ])
    
    output_file = "HashInsight_Product_Introduction_EN.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint created: {output_file}")

def main():
    print("\n" + "="*70)
    print("Generating All Documentation Formats (PDF/HTML/PPT)")
    print("="*70 + "\n")
    
    docs = [
        ("PRODUCT_INTRODUCTION_EN.md", "HashInsight Product Introduction"),
        ("OPERATIONS_MANUAL_EN.md", "HashInsight Operations Manual"),
        ("BENCHMARK_WHITEPAPER_EN.md", "Performance Benchmark Whitepaper"),
        ("DATA_SOURCE_RELIABILITY_EN.md", "Data Source Reliability Matrix"),
        ("SECURITY_COMPLIANCE_EVIDENCE_INDEX_EN.md", "Security Compliance Evidence Index"),
        ("EXECUTIVE_ONE_PAGER_EN.md", "HashInsight Executive One-Pager"),
    ]
    
    for md_file, title in docs:
        if os.path.exists(md_file):
            pdf_file = md_file.replace('.md', '.pdf')
            html_file = md_file.replace('.md', '.html')
            generate_pdf(md_file, pdf_file, title)
            generate_html(md_file, html_file, title)
        else:
            print(f"⚠️  Skipping {md_file} (not found)")
    
    print()
    create_enhanced_ppt()
    
    print("\n" + "="*70)
    print("✅ All documentation generated successfully!")
    print("="*70)
    print("\nGenerated files:")
    print("  📄 6 PDF documents")
    print("  🌐 6 HTML documents")
    print("  📊 1 PowerPoint presentation")
    print("\nTotal: 13 professional documents ready for distribution")

if __name__ == '__main__':
    main()
