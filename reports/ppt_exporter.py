"""
HashInsight Enterprise - Professional PPT Exporter
专业PPT导出器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
import logging
from datetime import datetime
from typing import List, Dict

logger = logging.getLogger(__name__)


class PPTExporter:
    """PPT报告导出器"""
    
    # 品牌颜色
    BRAND_COLORS = {
        'primary': RGBColor(13, 110, 253),  # Bootstrap Primary Blue
        'success': RGBColor(40, 167, 69),   # Green
        'danger': RGBColor(220, 53, 69),    # Red
        'dark': RGBColor(33, 37, 41),       # Dark Gray
        'light': RGBColor(248, 249, 250)    # Light Gray
    }
    
    def __init__(self, template_path: str = None):
        """初始化PPT导出器"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def create_title_slide(self, title: str, subtitle: str):
        """创建标题页"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"{subtitle}\n{datetime.now().strftime('%Y-%m-%d')}"
        
        # 设置字体样式
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.BRAND_COLORS['primary']
    
    def create_summary_slide(self, summary_data: Dict):
        """创建摘要页"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # 摘要内容
        content_top = Inches(1.5)
        for key, value in summary_data.items():
            text_box = slide.shapes.add_textbox(
                Inches(1), content_top, Inches(8), Inches(0.6)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"{key}: {value}"
            text_frame.paragraphs[0].font.size = Pt(18)
            content_top += Inches(0.8)
    
    def create_data_table(self, title: str, headers: List[str], data: List[List]):
        """创建数据表格页"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # 表格
        rows = len(data) + 1
        cols = len(headers)
        
        table = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5)
        ).table
        
        # 表头
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.BRAND_COLORS['primary']
        
        # 数据行
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
    
    def export(self, filename: str = None) -> bytes:
        """导出PPT"""
        if filename:
            self.prs.save(filename)
            return None
        else:
            output = io.BytesIO()
            self.prs.save(output)
            output.seek(0)
            return output.getvalue()
